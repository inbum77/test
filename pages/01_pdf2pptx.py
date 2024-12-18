import streamlit as st
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import io
from PIL import Image

# PDF를 이미지로 변환하는 함수
def pdf_to_images(pdf_bytes):
    pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
    images = []
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        pix = page.get_pixmap()
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        images.append(img)
    return images

# 이미지 리스트를 PPTX로 변환하는 함수
def images_to_pptx(images):
    prs = Presentation()
    
    for img in images:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 빈 슬라이드 레이아웃

        # 이미지의 크기를 인치로 변환
        width, height = img.size
        width_in_inches = width / 96
        height_in_inches = height / 96

        # 이미지 저장을 위해 BytesIO 사용
        image_stream = io.BytesIO()
        img.save(image_stream, format="PNG")
        image_stream.seek(0)

        # 슬라이드에 이미지 추가
        slide.shapes.add_picture(image_stream, Inches(0), Inches(0), width=Inches(width_in_inches), height=Inches(height_in_inches))

    return prs

# Streamlit UI
st.title("📄 PDF를 PPTX로 변환하기")
st.write("PDF 파일을 업로드하면 각 페이지를 이미지로 변환하여 PPTX 파일로 만들어줍니다. 😊")

# PDF 파일 업로드
uploaded_pdf = st.file_uploader("PDF 파일을 선택하세요:", type=["pdf"])

if uploaded_pdf:
    st.info("PDF 파일을 처리 중입니다... 잠시만 기다려주세요! ⏳")

    try:
        # PDF를 이미지로 변환
        pdf_bytes = uploaded_pdf.read()
        images = pdf_to_images(pdf_bytes)
        st.success(f"{len(images)} 페이지의 PDF가 성공적으로 변환되었습니다! 🖼️")

        # 이미지들을 PPTX로 변환
        pptx_presentation = images_to_pptx(images)

        # PPTX 파일을 메모리에서 저장
        pptx_bytes = io.BytesIO()
        pptx_presentation.save(pptx_bytes)
        pptx_bytes.seek(0)

        # 다운로드 버튼
        st.download_button(
            label="📥 PPTX 다운로드",
            data=pptx_bytes,
            file_name="converted_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

        st.success("PPTX 파일이 준비되었습니다! 다운로드 버튼을 클릭하세요. ✅")

    except Exception as e:
        st.error(f"오류가 발생했습니다: {e}")
